"""Draft 5분 발표 PPT 생성 — python-pptx.

deck.html 과 동일한 10 슬라이드. 16:9, 디자인 토큰 (cream/ink/muted) 적용.
발표자 노트는 script.md 의 해당 슬라이드 문단을 자동 삽입.

실행:
    python pitch/build-pptx.py
출력:
    pitch/deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pathlib import Path

# --- Design tokens (design.md 와 동일) ---
CREAM = RGBColor(0xFA, 0xF9, 0xF5)
INK = RGBColor(0x14, 0x14, 0x13)
MUTED = RGBColor(0x6C, 0x6A, 0x64)
MUTED_SOFT = RGBColor(0x8E, 0x8B, 0x82)
HAIRLINE = RGBColor(0xE6, 0xDF, 0xD8)
SURFACE_SOFT = RGBColor(0xF5, 0xF0, 0xE8)

FONT_BODY = "Pretendard"
FONT_DISPLAY = "함초롬바탕"

SLIDE_W_IN = 13.333  # 16:9 widescreen
SLIDE_H_IN = 7.5

OUT_PATH = Path(__file__).parent / "deck.pptx"
SCRIPT_PATH = Path(__file__).parent / "script.md"


# ─── Helpers ──────────────────────────────────────────────────────────────


def hex_rgb(hex_):
    h = hex_.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def set_bg(slide, rgb):
    """Cream 배경 — 슬라이드 가득 사각형."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb
    # 뒤로 보내기 — 다른 도형이 위에 보이게
    spTree = shape._element.getparent()
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def add_text(
    slide,
    text,
    left_in,
    top_in,
    width_in,
    height_in,
    *,
    font_name=FONT_BODY,
    font_size=22,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.4,
    letter_spacing=None,
):
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
    return tb


def add_eyebrow(slide, text, top_in=0.72):
    add_text(
        slide,
        text,
        left_in=0.83,
        top_in=top_in,
        width_in=12,
        height_in=0.4,
        font_size=11,
        color=MUTED,
        line_spacing=1.0,
    )


def add_page_num(slide, current, total):
    add_text(
        slide,
        f"{current} / {total}",
        left_in=11.5,
        top_in=7.0,
        width_in=1.5,
        height_in=0.3,
        font_size=10,
        color=MUTED_SOFT,
        align=PP_ALIGN.RIGHT,
    )


def add_divider(slide, top_in, *, left_in=0.83, width_in=11.67):
    """가는 헤어라인."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left_in),
        Inches(top_in),
        Inches(width_in),
        Emu(6350),  # ≈ 0.5 pt
    )
    line.line.fill.background()
    line.fill.solid()
    line.fill.fore_color.rgb = HAIRLINE


def add_notes(slide, notes_text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = notes_text


# ─── Slide builders ───────────────────────────────────────────────────────

NOTES = {
    1: "안녕하십니까. Draft 의 ○○○ 입니다. 저희는 한국 대학과 창업기관의 운영 매니저 한 분께만 가치를 전달하는 도구를 만들고 있습니다. 다음 5분 동안 이 회사가 풀고 있는 문제와, 왜 지금 이 시점에 풀려야 하는지 말씀드리겠습니다.",
    2: "저희가 인터뷰한 한 창업동아리 회장님은 매주 일요일 밤에 5시간을 씁니다. Drive 폴더 40개를 열어보고, 8개 팀이 무엇을 냈고 무엇을 빠뜨렸는지 시트에 옮겨 적습니다. 이게 매니저 한 명의 게으름 때문은 아닙니다. 학교가 매년 보고서를 요구하기 때문입니다.",
    3: "이 문제가 큰 이유는 돈이 모이고 있기 때문입니다. 2026년 서울 캠퍼스타운만 13개 대학에 312억이 들어갑니다. 서울 RISE 는 35개 대학에 5년간 4,225억입니다. LINC 3.0 은 연 3,025억입니다. 이 세 사업의 공통 조건이 정성적·정량적 활동 증빙입니다.",
    4: "Draft 는 한 문장으로 이렇게 설명됩니다. Drive 폴더만 연결하면, 마감과 요구사항과 결정사항이 자동으로 정리됩니다. 슬라이드를 잠깐 그대로 두고 청중이 읽게 함.",
    5: "화면을 보여드리겠습니다. 매니저가 본인의 Google Drive 폴더를 연결합니다. [전체 처리] 한 번 누르면, 안의 파일 12개를 저희가 차례로 다운로드해서 텍스트를 뽑아내고, 그 안에서 마감일, 요구사항, 외부 기관 이름, 결정사항을 항목 단위로 분리합니다. 매니저는 챗봇에게 한국어로 물으면 답이 빠르고 출처가 분명합니다.",
    6: "Notion 이나 Google Sheets 도 있습니다. 결정적 차이는 누가 입력하느냐입니다. 기존 도구는 학생이 형식에 맞춰 입력해야 작동합니다. 학생들은 안 합니다. Draft 는 학생에게 새 작업을 0건 추가합니다.",
    7: "지금 상황입니다. 6월 26일 FLIP 1기 8개 팀이 5주차 데모데이를 합니다. 이 운영을 저희 제품으로 진행합니다. 자체 검증에서 24개 항목 추출 정확도 100% 를 확인했습니다.",
    8: "누가 돈을 내는지 분명합니다. 1단계: 대학 창업지원단 직계약, 사업단 단위 SaaS, 1교당 연 600~1,200만 원. RISE 와 캠퍼스타운 예산의 콘텐츠·디지털 항목으로 합법적 집행. 2단계: 같은 학교 안에서 학과와 입주공간으로 확장.",
    9: "저는 FLIP 1기 운영자였습니다. 본인이 매주 일요일에 시트를 베끼던 사람입니다. 제 페인을 제가 풀고 있고, 첫 고객 8팀이 지금 저희 제품 안에서 일하고 있습니다.",
    10: "마지막입니다. 저희는 향후 12개월을 위해 Pre-seed 3~5억을 모집하고 있습니다. 정규직 엔지니어 1명과 영업 1명을 채용하고, 대학 5개 학교까지 PoC 를 확장합니다. KPI 는 유료 3교 계약, MAU 500 입니다. 들어주셔서 감사합니다.",
}


def build_cover(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "5분 발표 · FLIP 1기 데모데이", top_in=2.6)
    add_text(s, "Draft", 0.83, 3.0, 12, 1.6, font_name=FONT_DISPLAY, font_size=80, color=INK)
    add_text(
        s,
        "한국 창업기관 매니저의 운영 비서.",
        0.83, 4.7, 12, 0.6, font_size=22, color=MUTED, line_spacing=1.4,
    )
    add_text(s, "2026 · 발표자 ○○○", 0.83, 6.95, 6, 0.4, font_size=12, color=MUTED_SOFT)
    add_page_num(s, 1, total)
    add_notes(s, NOTES[1])


def build_problem(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "문제")
    add_text(
        s,
        "매니저 한 명이 매주 일요일 밤 5시간을\n시트에 베끼는 작업으로 보냅니다",
        0.83, 1.3, 12, 1.6, font_name=FONT_DISPLAY, font_size=44, line_spacing=1.2,
    )
    add_text(
        s,
        "FLIP 1기 회장 인터뷰 기준. 8개 팀, 5주차, Drive 폴더 40개, 카톡·이메일 3 채널. 이게 게으름이 아니라 학교 보고 양식이 그렇게 요구하기 때문입니다.",
        0.83, 3.7, 11.5, 1.5, font_size=20, color=hex_rgb("#3D3D3A"), line_spacing=1.6,
    )
    # 통계 3개
    stats = [("5h", "주당 운영 시간"), ("40+", "Drive 폴더"), ("3", "소통 채널")]
    for i, (num, label) in enumerate(stats):
        x = 0.83 + i * 4.1
        add_text(s, num, x, 5.4, 4, 1.0, font_name=FONT_DISPLAY, font_size=56, color=INK, line_spacing=1.0)
        add_text(s, label, x, 6.5, 4, 0.4, font_size=13, color=MUTED)
    add_page_num(s, 2, total)
    add_notes(s, NOTES[2])


def build_market(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "시장이 지금 신경 쓰는 이유")
    add_text(
        s,
        "학교가 매년 수천억을 쓰고,\n그 돈마다 활동 증빙이 따라붙습니다",
        0.83, 1.3, 12, 1.6, font_name=FONT_DISPLAY, font_size=40, line_spacing=1.2,
    )
    items = [
        ("서울 캠퍼스타운 2026", "312억 / 13개 대학"),
        ("서울 RISE 2025~2029", "4,225억 / 35개 대학"),
        ("LINC 3.0 연 사업비", "3,025억 / 년"),
        ("공통 요구사항", "정성·정량 활동 증빙"),
    ]
    for i, (name, amount) in enumerate(items):
        row = i // 2
        col = i % 2
        x = 0.83 + col * 6.0
        y = 3.6 + row * 1.6
        # 좌측 ink 바
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Emu(28575), Inches(1.3))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = INK
        add_text(s, name, x + 0.25, y + 0.05, 5.5, 0.4, font_size=14, color=MUTED)
        add_text(s, amount, x + 0.25, y + 0.5, 5.5, 0.7, font_name=FONT_DISPLAY, font_size=28, color=INK, line_spacing=1.1)
    add_text(s, "출처: 서울시 대학협력과 공시, 교육부 RISE 시행계획.", 0.83, 6.9, 8, 0.3, font_size=11, color=MUTED_SOFT)
    add_page_num(s, 3, total)
    add_notes(s, NOTES[3])


def build_solution(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "해결", top_in=2.0)
    add_text(
        s,
        "Drive 폴더만 연결하면,\n마감과 요구사항과 결정사항이",
        0.83, 2.6, 12, 2.0, font_name=FONT_DISPLAY, font_size=42, color=INK, line_spacing=1.35,
    )
    add_text(
        s,
        "자동으로 정리됩니다.",
        0.83, 4.2, 12, 1.0, font_name=FONT_DISPLAY, font_size=42, color=MUTED, line_spacing=1.35,
    )
    add_page_num(s, 4, total)
    add_notes(s, NOTES[4])


def build_demo(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "실제 작동")
    add_text(
        s,
        '"FLIP 1기" 폴더에서 추출된 항목',
        0.83, 1.3, 12, 0.7, font_name=FONT_DISPLAY, font_size=32, color=INK,
    )
    # frame
    frame = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.83), Inches(2.4), Inches(11.67), Inches(3.6)
    )
    frame.adjustments[0] = 0.04
    frame.line.color.rgb = HAIRLINE
    frame.line.width = Emu(6350)
    frame.fill.solid()
    frame.fill.fore_color.rgb = CREAM
    add_text(
        s, "Atom Digest · 마감 5 · 요구사항 7 · 결정 3",
        1.1, 2.55, 11, 0.35, font_size=12, color=MUTED,
    )
    cards = [
        ("마감", "3주차 사용성 테스트 결과 정리 완료", "3팀_2주차보고서.pdf · 2026-03-29"),
        ("요구사항", "60대 사용자 사용성 테스트 8건 진행", "3팀_2주차보고서.pdf · section 3"),
        ("결정", "LINC 중간평가 자료는 NRF 표준양식,\nPDF 10페이지 이내", "3팀_2주차보고서.pdf · section 3"),
    ]
    for i, (t, c, src) in enumerate(cards):
        x = 1.1 + i * 3.85
        y = 3.1
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.6), Inches(2.6))
        card.adjustments[0] = 0.07
        card.line.color.rgb = HAIRLINE
        card.line.width = Emu(6350)
        card.fill.solid()
        card.fill.fore_color.rgb = SURFACE_SOFT
        add_text(s, t, x + 0.2, y + 0.15, 3.3, 0.3, font_size=10, color=MUTED_SOFT)
        add_text(s, c, x + 0.2, y + 0.55, 3.3, 1.4, font_size=14, color=INK, line_spacing=1.5)
        add_text(s, src, x + 0.2, y + 2.05, 3.3, 0.5, font_size=10, color=MUTED_SOFT, font_name="Consolas")
    add_text(
        s,
        '매니저가 챗봇에게 "다가오는 마감 알려줘" 라고 묻습니다. 답이 출처 파일과 함께 즉시 나옵니다.',
        0.83, 6.4, 12, 0.5, font_size=18, color=hex_rgb("#3D3D3A"), line_spacing=1.5,
    )
    add_page_num(s, 5, total)
    add_notes(s, NOTES[5])


def build_diff(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "다른 점")
    add_text(
        s,
        "학생에게 새 작업을 0건 추가합니다",
        0.83, 1.3, 12, 0.8, font_name=FONT_DISPLAY, font_size=40, color=INK, line_spacing=1.2,
    )
    rows = [
        ("", "Notion / Sheets", "Draft"),
        ("누가 입력", "학생", "매니저만, 자동 추출"),
        ("형식", "형식 강제", "자유 형식 그대로"),
        ("활동 증빙", "매뉴얼 정리", "원본 출처 자동 인용"),
    ]
    col_widths = [3.0, 4.3, 4.3]
    col_xs = [0.83, 0.83 + col_widths[0], 0.83 + col_widths[0] + col_widths[1]]
    y_base = 3.0
    row_h = 0.8
    for ri, row in enumerate(rows):
        y = y_base + ri * row_h
        if ri == 0:
            for ci, cell in enumerate(row):
                add_text(s, cell, col_xs[ci] + 0.15, y + 0.15, col_widths[ci], 0.5, font_size=11, color=MUTED)
        else:
            for ci, cell in enumerate(row):
                if ci == 2:
                    pill = s.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(col_xs[ci]),
                        Inches(y + 0.05),
                        Inches(col_widths[ci] - 0.1),
                        Inches(row_h - 0.15),
                    )
                    pill.adjustments[0] = 0.2
                    pill.line.fill.background()
                    pill.fill.solid()
                    pill.fill.fore_color.rgb = INK
                    add_text(s, cell, col_xs[ci] + 0.2, y + 0.22, col_widths[ci], 0.5, font_size=16, color=CREAM, bold=True)
                else:
                    add_text(s, cell, col_xs[ci] + 0.15, y + 0.22, col_widths[ci], 0.5, font_size=16, color=INK)
        # row divider
        if ri > 0:
            add_divider(s, y, left_in=0.83, width_in=11.67)
    add_page_num(s, 6, total)
    add_notes(s, NOTES[6])


def build_traction(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "트랙션")
    add_text(s, "지금까지", 0.83, 1.3, 12, 1.0, font_name=FONT_DISPLAY, font_size=44, color=INK, line_spacing=1.2)
    stats = [
        ("8팀", "FLIP 1기 운영 중 (데모데이 6/26)"),
        ("100%", "사내 검증 24개 항목 정확도"),
        ("5h → 30분", "주당 운영 시간 단축 (예측)"),
    ]
    for i, (num, label) in enumerate(stats):
        x = 0.83 + i * 4.1
        add_text(s, num, x, 4.0, 4, 1.4, font_name=FONT_DISPLAY, font_size=48, color=INK, line_spacing=1.0)
        add_text(s, label, x, 5.6, 4, 1.0, font_size=14, color=MUTED, line_spacing=1.5)
    add_page_num(s, 7, total)
    add_notes(s, NOTES[7])


def build_bm(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "누가 돈을 내는가")
    add_text(
        s,
        "대학 창업지원단 직계약, 사업단 단위 SaaS",
        0.83, 1.3, 12, 0.8, font_name=FONT_DISPLAY, font_size=36, color=INK, line_spacing=1.2,
    )
    phases = [
        ("Phase 1 — 2026 하반기 ~ 2027",
         "대학당 연 600만 ~ 1,200만 원",
         "RISE·캠퍼스타운 콘텐츠·디지털 항목으로 합법적 집행. 의사결정자 = 창업지원단장 (부총장급 재량)."),
        ("Phase 2 — 2027 이후",
         "학교 안 확산 + 입주 공간",
         "같은 학교 안 학과·연구실 + 입주 기업 (TIPS, AC). 해외 진출은 QANDA 경로 — 학생 B2C 도구 이식."),
    ]
    for i, (label, title, desc) in enumerate(phases):
        x = 0.83 + i * 6.0
        y = 3.0
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(5.7), Inches(3.5))
        card.adjustments[0] = 0.04
        card.line.color.rgb = HAIRLINE
        card.line.width = Emu(6350)
        card.fill.solid()
        card.fill.fore_color.rgb = CREAM
        add_text(s, label, x + 0.3, y + 0.3, 5.2, 0.4, font_size=11, color=MUTED)
        add_text(s, title, x + 0.3, y + 0.75, 5.2, 0.8, font_name=FONT_DISPLAY, font_size=24, color=INK, line_spacing=1.2)
        add_text(s, desc, x + 0.3, y + 1.85, 5.2, 1.5, font_size=14, color=hex_rgb("#3D3D3A"), line_spacing=1.6)
    add_page_num(s, 8, total)
    add_notes(s, NOTES[8])


def build_team(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "팀")
    add_text(
        s,
        "본인이 매주 일요일에 시트를 베끼던 사람입니다",
        0.83, 1.3, 12, 1.0, font_name=FONT_DISPLAY, font_size=36, color=INK, line_spacing=1.2,
    )
    add_text(
        s,
        "Founder · CEO  ○○○\nFLIP 1기 운영. 매니저로서 본인의 페인을 본인이 풀고 있습니다.\n첫 고객 8팀이 지금 저희 제품 안에서 운영되고 있습니다.",
        0.83, 3.0, 11.5, 2.5, font_size=18, color=INK, line_spacing=1.7,
    )
    add_text(
        s,
        "제품 엔지니어링 (Claude Code 협업) — Next.js 15 + Supabase + 멀티 LLM (Claude · Gemini fallback).",
        0.83, 5.5, 11.5, 1.0, font_size=15, color=MUTED, line_spacing=1.6,
    )
    add_page_num(s, 9, total)
    add_notes(s, NOTES[9])


def build_ask(s, total):
    set_bg(s, CREAM)
    add_eyebrow(s, "다음 12개월")
    add_text(
        s,
        "Pre-seed 3 ~ 5 억 모집 중",
        0.83, 1.3, 12, 1.0, font_name=FONT_DISPLAY, font_size=44, color=INK, line_spacing=1.2,
    )
    items = [
        "·  정규직 엔지니어 1인 + 영업 1인 채용",
        "·  대학 5개 학교로 PoC 확장",
        "·  KPI: 유료 3교 계약 · MAU 500",
    ]
    for i, it in enumerate(items):
        add_text(s, it, 0.83, 3.2 + i * 0.7, 12, 0.6, font_size=22, color=INK, line_spacing=1.5)
    add_text(
        s,
        "오늘 자리에서 한 분이라도 follow-up 미팅 가능하시면, 끝나고 명함 부탁드립니다.",
        0.83, 5.8, 11.5, 0.8, font_size=18, color=hex_rgb("#3D3D3A"), line_spacing=1.6,
    )
    add_text(
        s, "○○○ · email@draft.kr · 010-XXXX-XXXX",
        0.83, 6.9, 11, 0.4, font_size=12, color=MUTED_SOFT,
    )
    add_page_num(s, 10, total)
    add_notes(s, NOTES[10])


# ─── Main ─────────────────────────────────────────────────────────────────


def main():
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    builders = [
        build_cover, build_problem, build_market, build_solution,
        build_demo, build_diff, build_traction, build_bm,
        build_team, build_ask,
    ]
    total = len(builders)
    for build in builders:
        slide = prs.slides.add_slide(blank)
        build(slide, total)

    prs.save(str(OUT_PATH))
    print(f"OK → {OUT_PATH} ({OUT_PATH.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
